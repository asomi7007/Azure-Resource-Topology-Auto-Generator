from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from .icon_manager import IconManager
import os
import tempfile

def generate_pptx_file(layout_nodes, relationships, output_path):
    icon_mgr = IconManager()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank Layout
    
    shape_map = {}
    
    # 1. Draw Nodes (Icons)
    for node in layout_nodes:
        x_cm = Cm(node['x'] / 30.0)
        y_cm = Cm(node['y'] / 30.0)
        w_cm = Cm(node['w'] / 30.0)
        
        res_type = node['resource']['type'].lower()
        
        # Get Icon
        icon_img = icon_mgr.get_icon_image(res_type, 64, 64)
        
        shape = None
        if icon_img:
            # Save dump
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                icon_img.save(tmp.name)
                tmp_path = tmp.name
            
            # Add Picture
            icon_size = Cm(1.5)
            # Center in the layout box
            ix = x_cm + (w_cm - icon_size)/2
            iy = y_cm
            
            shape = slide.shapes.add_picture(tmp_path, ix, iy, width=icon_size, height=icon_size)
            os.remove(tmp_path)
            
        else:
            # Fallback
            shape = slide.shapes.add_shape(5, x_cm, y_cm, Cm(1.5), Cm(1.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(200, 200, 200)

        # Map ID -> Shape (for connectors)
        if shape:
            shape_map[node['id'].lower()] = shape
            
        # Text Label
        tx = x_cm
        ty = y_cm + Cm(1.6)
        label = slide.shapes.add_textbox(tx, ty, w_cm, Cm(1))
        label.text_frame.text = node['resource']['name']
        label.text_frame.paragraphs[0].font.size = Pt(9)

    # 2. Draw Connectors
    for rel in relationships:
        src_id = rel['from'].lower()
        dst_id = rel['to'].lower()
        
        if src_id in shape_map and dst_id in shape_map:
            src_shape = shape_map[src_id]
            dst_shape = shape_map[dst_id]
            
            # Elbow Connector
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.ELBOW, 0, 0, 0, 0
            )
            
            # Connect closest points (Auto)
            # Or force specific sides?
            # VNet (Left) -> Subnet (Right)
            # Graph flows Left -> Right
            # Let's let PPT handle auto route
            connector.begin_connect(src_shape, 3) # 3=Right?
            connector.end_connect(dst_shape, 1)   # 1=Left?
            
            # Style
            line = connector.line
            category = rel.get('category')
            
            if category == 'Traffic':
                 line.color.rgb = RGBColor(0, 180, 0)
                 line.width = Pt(2)
                 line.head_end = 2 # Arrow
            elif category == 'Physical':
                 # VNet -> Subnet
                 line.color.rgb = RGBColor(0, 120, 212)
                 line.width = Pt(1.5)
            else:
                 line.color.rgb = RGBColor(120, 120, 120)
                 line.width = Pt(1)
                 line.dash_style = 4 # SquareDot

    prs.save(output_path)
